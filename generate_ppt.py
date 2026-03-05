from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colour Palette (matches CampusHub dark theme) ──────────────────────────
BG        = RGBColor(0x09, 0x09, 0x0F)   # near-black background
SURFACE   = RGBColor(0x11, 0x11, 0x1B)   # dark surface
SURFACE2  = RGBColor(0x1A, 0x1A, 0x2E)   # card surface
PRIMARY   = RGBColor(0x6C, 0x63, 0xFF)   # electric purple
CYAN      = RGBColor(0x00, 0xD4, 0xFF)   # cyan accent
AMBER     = RGBColor(0xFF, 0x9F, 0x43)   # amber (Lost & Found)
GREEN     = RGBColor(0x00, 0xB8, 0x94)   # green (Textbooks)
ROSE      = RGBColor(0xFF, 0x76, 0x75)   # rose (warning/lost)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
TEXT      = RGBColor(0xE2, 0xE8, 0xF0)
MUTED     = RGBColor(0x8B, 0x8B, 0xA7)
PURPLE2   = RGBColor(0xA7, 0x8B, 0xFA)

W = Inches(13.33)   # widescreen 16:9
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]   # completely blank

# ── Helpers ────────────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=None, alpha=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)   # MSO_SHAPE_TYPE.RECTANGLE
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        if line_w: shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, x, y, w, h,
             size=18, bold=False, color=TEXT, align=PP_ALIGN.LEFT,
             italic=False, wrap=True, name="Inter"):
    txb = slide.shapes.add_textbox(x, y, w, h)
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    run.font.name  = name
    return txb

def add_multiline(slide, lines, x, y, w, h, size=14, color=MUTED, spacing=Pt(6), name="Inter"):
    """lines = list of (text, bold, color) tuples"""
    txb = slide.shapes.add_textbox(x, y, w, h)
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for (txt, bold, col) in lines:
        if first:
            p = tf.paragraphs[0]; first = False
        else:
            p = tf.add_paragraph()
        p.space_before = spacing
        run = p.add_run()
        run.text = txt
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = col
        run.font.name = name

def slide_bg(slide, color=BG):
    bg = add_rect(slide, 0, 0, W, H, fill=color)
    bg.name = "Background"

def glow_orb(slide, cx, cy, size, color, opacity_hint=30):
    """Approximate glow with layered semi-transparent circles."""
    for factor in [1.0, 0.65, 0.35]:
        s = Inches(size * factor)
        shape = slide.shapes.add_shape(9, cx - s/2, cy - s/2, s, s)  # oval
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        # No true transparency in pptx without XML hacking; use lighter tones
    # Tiny bright core
    core = Inches(size * 0.12)
    shape2 = slide.shapes.add_shape(9, cx - core/2, cy - core/2, core, core)
    shape2.fill.solid()
    shape2.fill.fore_color.rgb = WHITE
    shape2.line.fill.background()

def accent_bar(slide, color=PRIMARY):
    add_rect(slide, 0, H - Inches(0.08), W, Inches(0.08), fill=color)

def slide_number(slide, num, total=12):
    add_text(slide, f"{num} / {total}", W - Inches(1.2), H - Inches(0.45),
             Inches(1), Inches(0.35), size=10, color=MUTED, align=PP_ALIGN.RIGHT)

def tag_pill(slide, label, x, y, color=PRIMARY, bg_alpha=None):
    w = Inches(len(label)*0.11 + 0.4)
    h = Inches(0.32)
    box = add_rect(slide, x, y, w, h, fill=SURFACE2, line=color, line_w=Pt(0.8))
    add_text(slide, label, x + Inches(0.1), y + Inches(0.04), w, h,
             size=10, bold=True, color=color, align=PP_ALIGN.LEFT)
    return w

def section_icon(slide, emoji, x, y, size=Inches(0.65), bg=SURFACE2, fg=PRIMARY):
    box = add_rect(slide, x, y, size, size, fill=bg, line=fg, line_w=Pt(0.6))
    box.adjustments[0] = 0.1   # no-op, just styling hint
    add_text(slide, emoji, x + Inches(0.08), y + Inches(0.04),
             size, size, size=22, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
slide_bg(sl)

# Orbs
glow_orb(sl, Inches(1.5),  Inches(1.5),  5.5, PRIMARY)
glow_orb(sl, Inches(11.5), Inches(5.5),  4.0, CYAN)
glow_orb(sl, Inches(7.0),  Inches(6.5),  2.5, AMBER)

# Grid overlay (thin line box pattern suggested via subtle rect)
add_rect(sl, 0, 0, W, H, fill=None, line=None)

# Badge
tag_pill(sl, " 🎓  MERN Stack  |  Full-Stack Web App ", Inches(4.2), Inches(1.4), color=CYAN)

# Title
add_text(sl, "CampusHub", Inches(2.0), Inches(2.0), Inches(9.33), Inches(1.5),
         size=68, bold=True, color=WHITE, align=PP_ALIGN.CENTER, name="Calibri")

# Subtitle gradient hint (two coloured spans not supported; use single)
add_text(sl, "Digital Students Hub", Inches(2.5), Inches(3.35), Inches(8.33), Inches(0.7),
         size=28, bold=False, color=CYAN, align=PP_ALIGN.CENTER)

add_text(sl,
         "Campus News  ·  Lost & Found  ·  Used Textbook Marketplace",
         Inches(1.5), Inches(4.05), Inches(10.33), Inches(0.5),
         size=15, color=MUTED, align=PP_ALIGN.CENTER)

# Tech pills row
px = Inches(3.4);  gap = Inches(1.55)
for label, col in [("MongoDB", GREEN), ("Express", MUTED), ("React", CYAN), ("Node.js", GREEN)]:
    tag_pill(sl, f"  {label}  ", px, Inches(5.1), color=col)
    px += gap

accent_bar(sl)
slide_number(sl, 1)

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — PROJECT OVERVIEW
# ══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
slide_bg(sl)
glow_orb(sl, Inches(12), Inches(1), 3, PRIMARY)

add_text(sl, "Project Overview", Inches(0.6), Inches(0.55), Inches(8), Inches(0.6),
         size=34, bold=True, color=WHITE, name="Calibri")
add_rect(sl, Inches(0.6), Inches(1.2), Inches(1.2), Inches(0.04), fill=PRIMARY)

add_text(sl,
    "CampusHub is a full-stack MERN web application that acts as a central digital "
    "hub for university students — enabling them to share campus news, post lost & found "
    "items, and buy or sell used textbooks, all in one beautifully designed platform.",
    Inches(0.6), Inches(1.4), Inches(8.5), Inches(1.2),
    size=14, color=TEXT)

# 3 stat boxes
for i, (num, lbl, col) in enumerate([
        ("3", "Core Modules", PRIMARY),
        ("9+", "API Endpoints", CYAN),
        ("∞", "Students Helped", GREEN)]):
    bx = Inches(0.6 + i * 4.15)
    add_rect(sl, bx, Inches(2.75), Inches(3.8), Inches(1.3), fill=SURFACE2, line=col, line_w=Pt(1.5))
    add_text(sl, num, bx + Inches(0.2), Inches(2.85), Inches(3.4), Inches(0.65),
             size=42, bold=True, color=col, align=PP_ALIGN.CENTER, name="Calibri")
    add_text(sl, lbl, bx + Inches(0.2), Inches(3.5), Inches(3.4), Inches(0.4),
             size=13, color=MUTED, align=PP_ALIGN.CENTER)

# Feature bullets
bullets = [
    ("📰", "Campus News",     "Post & browse announcements, events, academic & sports updates", PRIMARY),
    ("🔍", "Lost & Found",    "Report lost items, post found items — filter by category & location", AMBER),
    ("📚", "Used Textbooks",  "List & discover second-hand books — filter by subject, condition & price", GREEN),
    ("📞", "Direct Contact",  "Every listing shows poster's email & phone — no login required", CYAN),
]
for i, (icon, title, desc, col) in enumerate(bullets):
    by = Inches(4.25) + i * Inches(0.72)
    add_rect(sl, Inches(0.6), by, Inches(12.1), Inches(0.62), fill=SURFACE, line=col, line_w=Pt(0.6))
    add_text(sl, icon, Inches(0.75), by + Inches(0.1), Inches(0.5), Inches(0.45), size=16)
    add_text(sl, title, Inches(1.35), by + Inches(0.1), Inches(2.0), Inches(0.45), size=13, bold=True, color=col)
    add_text(sl, desc,  Inches(3.4),  by + Inches(0.1), Inches(9.1), Inches(0.45), size=12, color=MUTED)

accent_bar(sl)
slide_number(sl, 2)

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — TECH STACK
# ══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
slide_bg(sl)
glow_orb(sl, Inches(0.5), Inches(6), 3.5, CYAN)

add_text(sl, "Tech Stack — MERN", Inches(0.6), Inches(0.55), Inches(10), Inches(0.6),
         size=34, bold=True, color=WHITE, name="Calibri")
add_rect(sl, Inches(0.6), Inches(1.2), Inches(1.4), Inches(0.04), fill=CYAN)

cards = [
    ("M", "MongoDB", "NoSQL Database", GREEN,
     ["Mongoose ODM", "3 data models", "Atlas / Local", "Schema validation"]),
    ("E", "Express.js", "REST API Backend", MUTED,
     ["RESTful routes", "CORS enabled", "Modular routing", "Error handling"]),
    ("R", "React 18", "Frontend UI", CYAN,
     ["Vite build tool", "React Router v6", "Hooks & Context", "Lucide icons"]),
    ("N", "Node.js", "Runtime Environment", GREEN,
     ["v20+ LTS", "dotenv config", "nodemon (dev)", "npm workspaces"]),
]

for i, (letter, name, role, col, points) in enumerate(cards):
    cx = Inches(0.55 + i * 3.22)
    cy = Inches(1.5)
    cw = Inches(3.05)
    ch = Inches(5.6)
    add_rect(sl, cx, cy, cw, ch, fill=SURFACE2, line=col, line_w=Pt(1.5))
    # Letter badge
    add_rect(sl, cx + Inches(0.95), cy + Inches(0.2), Inches(1.15), Inches(1.0), fill=col)
    add_text(sl, letter, cx + Inches(0.95), cy + Inches(0.18), Inches(1.15), Inches(1.0),
             size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER, name="Calibri")
    add_text(sl, name, cx + Inches(0.1), cy + Inches(1.3), cw - Inches(0.2), Inches(0.45),
             size=18, bold=True, color=col, align=PP_ALIGN.CENTER, name="Calibri")
    add_text(sl, role, cx + Inches(0.1), cy + Inches(1.72), cw - Inches(0.2), Inches(0.35),
             size=11, color=MUTED, align=PP_ALIGN.CENTER)
    add_rect(sl, cx + Inches(0.3), cy + Inches(2.15), cw - Inches(0.6), Inches(0.02), fill=SURFACE)
    for j, pt in enumerate(points):
        py = cy + Inches(2.3) + j * Inches(0.68)
        add_rect(sl, cx + Inches(0.25), py, Inches(0.06), Inches(0.06),
                 fill=col)
        add_text(sl, pt, cx + Inches(0.42), py - Inches(0.06), cw - Inches(0.6), Inches(0.5),
                 size=12, color=TEXT)

accent_bar(sl)
slide_number(sl, 3)

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — SYSTEM ARCHITECTURE
# ══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
slide_bg(sl)

add_text(sl, "System Architecture", Inches(0.6), Inches(0.5), Inches(10), Inches(0.6),
         size=34, bold=True, color=WHITE, name="Calibri")
add_rect(sl, Inches(0.6), Inches(1.15), Inches(1.6), Inches(0.04), fill=PRIMARY)

# Architecture diagram — three boxes with arrows
def arch_box(slide, x, y, w, h, title, subtitle, items, col):
    add_rect(slide, x, y, w, h, fill=SURFACE2, line=col, line_w=Pt(1.5))
    add_rect(slide, x, y, w, Inches(0.55), fill=col)
    add_text(slide, title,    x + Inches(0.15), y + Inches(0.07), w - Inches(0.3), Inches(0.38),
             size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER, name="Calibri")
    add_text(slide, subtitle, x + Inches(0.1),  y + Inches(0.55), w - Inches(0.2), Inches(0.35),
             size=11, color=col, align=PP_ALIGN.CENTER)
    for j, it in enumerate(items):
        add_text(slide, f"▸  {it}", x + Inches(0.2), y + Inches(0.98) + j*Inches(0.47),
                 w - Inches(0.3), Inches(0.4), size=11, color=MUTED)

bw = Inches(3.55)
bh = Inches(4.8)
by = Inches(1.45)
gap = Inches(0.88)

arch_box(sl, Inches(0.5), by, bw, bh, "React Frontend", "localhost:3000",
         ["React 18 + Vite", "React Router v6", "Lucide Icons", "react-hot-toast", "Axios HTTP client", "CSS Variables"], CYAN)

arch_box(sl, Inches(0.5) + bw + gap, by, bw, bh, "Express Backend", "localhost:5000",
         ["server.js (entry)", "/api/news", "/api/lostfound", "/api/textbooks", "/api/stats", "CORS + dotenv"], PRIMARY)

arch_box(sl, Inches(0.5) + (bw + gap)*2, by, bw, bh, "MongoDB Database", "mongoose ODM",
         ["News model", "LostFound model", "Textbook model", "Schema validation", "Atlas / Local", "Timestamps"], GREEN)

# Arrows
for ax in [Inches(0.5) + bw + Inches(0.04), Inches(0.5) + bw*2 + gap + Inches(0.04)]:
    ay = by + bh/2 - Inches(0.08)
    add_rect(sl, ax, ay, gap - Inches(0.08), Inches(0.07), fill=SURFACE)
    # arrowhead suggestion (triangle approximation)
    add_text(sl, "⟶", ax, ay - Inches(0.18), gap, Inches(0.4), size=22, color=MUTED, align=PP_ALIGN.CENTER)

# Vite proxy note
add_text(sl, "⚡ Vite Proxy: /api → :5000 (dev) · Build: static bundle served by Node in prod",
         Inches(0.5), Inches(6.5), Inches(12.33), Inches(0.4),
         size=11, color=MUTED, italic=True, align=PP_ALIGN.CENTER)

accent_bar(sl)
slide_number(sl, 4)

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — DATABASE MODELS
# ══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
slide_bg(sl)
glow_orb(sl, Inches(12.5), Inches(6.5), 4, GREEN)

add_text(sl, "MongoDB Data Models", Inches(0.6), Inches(0.5), Inches(10), Inches(0.6),
         size=34, bold=True, color=WHITE, name="Calibri")
add_rect(sl, Inches(0.6), Inches(1.15), Inches(1.5), Inches(0.04), fill=GREEN)

def model_card(slide, x, y, w, title, col, fields):
    h = Inches(0.52 + len(fields) * 0.5)
    add_rect(slide, x, y, w, h, fill=SURFACE2, line=col, line_w=Pt(1.2))
    add_rect(slide, x, y, w, Inches(0.48), fill=col)
    add_text(slide, title, x + Inches(0.15), y + Inches(0.07), w - Inches(0.3), Inches(0.36),
             size=14, bold=True, color=WHITE, name="Calibri")
    for i, (fname, ftype) in enumerate(fields):
        fy = y + Inches(0.55) + i * Inches(0.48)
        add_text(slide, fname, x + Inches(0.2),  fy, Inches(1.6), Inches(0.42), size=11, bold=True, color=TEXT)
        add_text(slide, ftype, x + Inches(1.85), fy, w - Inches(2.1), Inches(0.42), size=11, color=MUTED)

mw = Inches(4.0)
model_card(sl, Inches(0.5), Inches(1.4), mw, "📰  News",  PRIMARY,
    [("title",       "String  ·  required"),
     ("content",     "String  ·  required"),
     ("category",    "Enum: announcement / event / academic / sports / cultural / other"),
     ("authorName",  "String  ·  required"),
     ("authorEmail", "String  ·  optional"),
     ("authorPhone", "String  ·  optional"),
     ("tags",        "[ String ]"),
     ("timestamps",  "createdAt  ·  updatedAt")])

model_card(sl, Inches(4.72), Inches(1.4), mw, "🔍  LostFound", AMBER,
    [("type",         "Enum: lost | found  ·  required"),
     ("itemName",     "String  ·  required"),
     ("category",     "Enum: electronics / keys / bags / books / clothing / …"),
     ("description",  "String  ·  required"),
     ("location",     "String  ·  required"),
     ("dateLostFound","Date  ·  required"),
     ("posterName",   "String  ·  required"),
     ("status",       "Enum: active | resolved  ·  default: active")])

model_card(sl, Inches(8.94), Inches(1.4), mw, "📚  Textbook", GREEN,
    [("title",       "String  ·  required"),
     ("author",      "String  ·  required"),
     ("subject",     "String  ·  required"),
     ("category",    "Enum: engineering / science / cs / arts / business / …"),
     ("condition",   "Enum: like-new | good | fair | poor  ·  required"),
     ("price",       "Number  ·  required  ·  min 0"),
     ("negotiable",  "Boolean  ·  default true"),
     ("sold",        "Boolean  ·  default false")])

accent_bar(sl)
slide_number(sl, 5)

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 6 — API ENDPOINTS
# ══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
slide_bg(sl)

add_text(sl, "REST API Endpoints", Inches(0.6), Inches(0.5), Inches(10), Inches(0.6),
         size=34, bold=True, color=WHITE, name="Calibri")
add_rect(sl, Inches(0.6), Inches(1.15), Inches(1.5), Inches(0.04), fill=CYAN)

def method_pill(slide, x, y, method):
    col = {"GET": GREEN, "POST": PRIMARY, "PATCH": AMBER, "DELETE": ROSE}[method]
    w = Inches(0.75)
    add_rect(slide, x, y, w, Inches(0.36), fill=col)
    add_text(slide, method, x, y + Inches(0.04), w, Inches(0.3), size=10, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER)
    return w

endpoints = [
    # (method, path, description)
    ("GET",    "/api/news",                   "Fetch all news  ·  filters: category, search, page"),
    ("POST",   "/api/news",                   "Create a news post"),
    ("DELETE", "/api/news/:id",               "Delete a news post"),
    ("GET",    "/api/lostfound",              "Fetch all items  ·  filters: type, category, status, search"),
    ("POST",   "/api/lostfound",              "Create a lost/found report"),
    ("PATCH",  "/api/lostfound/:id/status",   "Update status → 'resolved'"),
    ("DELETE", "/api/lostfound/:id",          "Delete a lost/found item"),
    ("GET",    "/api/textbooks",              "Fetch books  ·  filters: category, condition, price, sort"),
    ("POST",   "/api/textbooks",             "Create a textbook listing"),
    ("PATCH",  "/api/textbooks/:id/sold",     "Mark textbook as sold"),
    ("DELETE", "/api/textbooks/:id",          "Delete a textbook listing"),
    ("GET",    "/api/stats",                  "Dashboard stats: news, lost, found, textbook counts"),
]

for i, (method, path, desc) in enumerate(endpoints):
    ry = Inches(1.45) + i * Inches(0.49)
    row_col = SURFACE2 if i % 2 == 0 else SURFACE
    add_rect(sl, Inches(0.5), ry, Inches(12.33), Inches(0.45), fill=row_col)
    method_pill(sl, Inches(0.55), ry + Inches(0.04), method)
    add_text(sl, path, Inches(1.4),  ry + Inches(0.06), Inches(3.5),  Inches(0.38), size=11, bold=True, color=TEXT, name="Courier New")
    add_text(sl, desc, Inches(5.0),  ry + Inches(0.06), Inches(7.7),  Inches(0.38), size=11, color=MUTED)

accent_bar(sl)
slide_number(sl, 6)

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 7 — CAMPUS NEWS FEATURE
# ══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
slide_bg(sl)
glow_orb(sl, Inches(11), Inches(1.5), 4, PRIMARY)

add_text(sl, "📰  Campus News", Inches(0.6), Inches(0.5), Inches(10), Inches(0.6),
         size=34, bold=True, color=WHITE, name="Calibri")
add_rect(sl, Inches(0.6), Inches(1.15), Inches(1.0), Inches(0.04), fill=PRIMARY)

# Feature list
features = [
    ("Post News",       "Students can share announcements, event notices, academic updates, sports results, and cultural events."),
    ("6 Categories",    "Announcement · Event · Academic · Sports · Cultural · Other — each visually colour-coded."),
    ("Smart Search",    "Full-text regex search across title, content, and tags simultaneously on the backend."),
    ("Tag System",      "Free-form tags (comma-separated) stored as array — allows custom topic keywording."),
    ("Author Contact",  "Author name, email and phone displayed on every card — click email to open mail client."),
    ("Delete Posts",    "One-click delete with confirmation — instantly synced with MongoDB via REST DELETE."),
]
for i, (title, desc) in enumerate(features):
    fx = Inches(0.6) if i < 3 else Inches(6.77)
    fy = Inches(1.5) + (i % 3) * Inches(1.7)
    add_rect(sl, fx, fy, Inches(5.7), Inches(1.55), fill=SURFACE2, line=PRIMARY, line_w=Pt(0.8))
    add_text(sl, f"✦  {title}", fx + Inches(0.2), fy + Inches(0.18), Inches(5.3), Inches(0.4),
             size=14, bold=True, color=PRIMARY)
    add_text(sl, desc, fx + Inches(0.2), fy + Inches(0.58), Inches(5.3), Inches(0.85),
             size=11, color=MUTED)

# Mock card preview
cx = Inches(0.6)
cy = Inches(6.2)
add_rect(sl, cx, cy, Inches(12.2), Inches(0.75), fill=SURFACE, line=SURFACE2)
add_text(sl, "📋  Sample Card :", cx + Inches(0.2), cy + Inches(0.18), Inches(2.2), Inches(0.4),
         size=11, color=MUTED)
for label, col, offset in [("Event", CYAN, 0), ("Tech Fest 2026", WHITE, 0.75), ("📅 2d ago", MUTED, 4.2), ("✉ author@univ.edu", PRIMARY, 5.8)]:
    add_text(sl, label, cx + Inches(offset + 2.35), cy + Inches(0.2), Inches(2.0), Inches(0.38),
             size=11, bold=(col == WHITE), color=col)

accent_bar(sl)
slide_number(sl, 7)

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 8 — LOST & FOUND FEATURE
# ══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
slide_bg(sl)
glow_orb(sl, Inches(11), Inches(6), 4, AMBER)

add_text(sl, "🔍  Lost & Found", Inches(0.6), Inches(0.5), Inches(10), Inches(0.6),
         size=34, bold=True, color=WHITE, name="Calibri")
add_rect(sl, Inches(0.6), Inches(1.15), Inches(1.1), Inches(0.04), fill=AMBER)

# Left column
left_features = [
    ("Type Toggle",       ROSE,  "Switch between 🔍 Lost and ✅ Found — colour-coded cards instantly update."),
    ("10 Item Categories",AMBER, "Electronics · ID/Cards · Keys · Bags · Books · Clothing · Jewelry · Sports · Stationery · Other"),
    ("Location & Date",   CYAN,  "Every report includes where and when the item was lost/found."),
    ("Mark as Resolved",  GREEN, "Poster can mark their item as resolved once reunited — dims the card visually."),
    ("Direct Contact",    PRIMARY,"Poster's email and phone shown on every card for direct student-to-student connection."),
]
for i, (title, col, desc) in enumerate(left_features):
    fy = Inches(1.45) + i * Inches(1.05)
    add_rect(sl, Inches(0.5), fy, Inches(7.8), Inches(0.95), fill=SURFACE2, line=col, line_w=Pt(1.0))
    add_text(sl, f"▸ {title}", Inches(0.7), fy + Inches(0.1), Inches(3.5), Inches(0.4),
             size=13, bold=True, color=col)
    add_text(sl, desc, Inches(0.7), fy + Inches(0.48), Inches(7.5), Inches(0.42),
             size=11, color=MUTED)

# Right column — workflow box
add_rect(sl, Inches(8.5), Inches(1.45), Inches(4.35), Inches(5.6), fill=SURFACE2, line=AMBER, line_w=Pt(1.2))
add_text(sl, "Student Workflow", Inches(8.7), Inches(1.6), Inches(4.0), Inches(0.4),
         size=14, bold=True, color=AMBER, align=PP_ALIGN.CENTER)

steps = [
    ("1️⃣", "Click 'Post Item'", WHITE),
    ("",    "Select Lost or Found", MUTED),
    ("2️⃣", "Fill in item details", WHITE),
    ("",    "Name, category, location, date", MUTED),
    ("3️⃣", "Add contact info", WHITE),
    ("",    "Email and / or phone number", MUTED),
    ("4️⃣", "Browse the board", WHITE),
    ("",    "Filter, search, find a match", MUTED),
    ("5️⃣", "Contact directly", WHITE),
    ("",    "Email or call the poster", MUTED),
    ("✅",  "Mark Resolved", GREEN),
]
for i, (icon, text, col) in enumerate(steps):
    sy = Inches(2.1) + i * Inches(0.43)
    add_text(sl, f"{icon}  {text}", Inches(8.7), sy, Inches(4.0), Inches(0.38),
             size=11, color=col, bold=(col == WHITE))

accent_bar(sl)
slide_number(sl, 8)

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 9 — USED TEXTBOOKS FEATURE
# ══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
slide_bg(sl)
glow_orb(sl, Inches(0.5), Inches(1), 3.5, GREEN)

add_text(sl, "📚  Used Textbook Marketplace", Inches(0.6), Inches(0.5), Inches(12), Inches(0.6),
         size=34, bold=True, color=WHITE, name="Calibri")
add_rect(sl, Inches(0.6), Inches(1.15), Inches(1.8), Inches(0.04), fill=GREEN)

# Feature table
headers = ["Feature", "Details"]
rows = [
    ("10 Subject Categories",   "Engineering · Science · Maths · CS · Arts · Business · Medicine · Law · Social Science"),
    ("4 Condition Levels",      "Like New 🟢 · Good 🟣 · Fair 🟠 · Poor 🔴  —  colour-coded badges on every card"),
    ("Price Filtering",         "Min/max price range filter + Sort: Newest / Price Low→High / Price High→Low"),
    ("Negotiable Flag",         "Checkbox on listing form — shown as a cyan 'Negotiable' badge on the card"),
    ("Contact Seller Modal",    "Clicking 'Contact' opens a styled modal with seller email (mailto link) & phone"),
    ("Mark as Sold",            "PATCH /api/textbooks/:id/sold — card shows a SOLD overlay with dim effect"),
    ("Search",                  "Full-text search across title, author, and subject fields via MongoDB regex"),
    ("Pagination-Ready API",    "Backend supports ?page=&limit= params — frontend ready for paginator integration"),
]

add_rect(sl, Inches(0.5), Inches(1.45), Inches(12.33), Inches(0.46), fill=GREEN)
add_text(sl, headers[0], Inches(0.7), Inches(1.5), Inches(3.0), Inches(0.38), size=12, bold=True, color=WHITE)
add_text(sl, headers[1], Inches(3.8), Inches(1.5), Inches(9.0), Inches(0.38), size=12, bold=True, color=WHITE)

for i, (feat, detail) in enumerate(rows):
    ry = Inches(1.91) + i * Inches(0.64)
    add_rect(sl, Inches(0.5), ry, Inches(12.33), Inches(0.6), fill=SURFACE2 if i%2==0 else SURFACE)
    add_text(sl, feat,   Inches(0.65), ry + Inches(0.1), Inches(3.1), Inches(0.44), size=12, bold=True, color=WHITE)
    add_text(sl, detail, Inches(3.8),  ry + Inches(0.1), Inches(8.8), Inches(0.44), size=11, color=MUTED)

accent_bar(sl)
slide_number(sl, 9)

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 10 — UI / UX DESIGN HIGHLIGHTS
# ══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
slide_bg(sl)
glow_orb(sl, Inches(6.5), Inches(0.5), 5, PRIMARY)

add_text(sl, "UI / UX Design Highlights", Inches(0.6), Inches(0.5), Inches(10), Inches(0.6),
         size=34, bold=True, color=WHITE, name="Calibri")
add_rect(sl, Inches(0.6), Inches(1.15), Inches(1.7), Inches(0.04), fill=PURPLE2)

design_points = [
    ("🌑", "Dark Glassmorphism",   PRIMARY,  "Near-black background (#09090F) with blurred backdrop-filter navbar and glowing orbs."),
    ("🎨", "CSS Variables",         CYAN,     "Full design token system — colours, radii, shadows, gradients defined as :root vars."),
    ("✨", "Micro-Animations",      PURPLE2,  "fade-in cards, hover translateY lift, shimmer skeleton loaders, spinning orbs on hero."),
    ("📱", "Responsive Layout",     GREEN,    "CSS Grid auto-fill columns collapse to 1-col on mobile. Hamburger nav on ≤768px."),
    ("🔔", "Toast Notifications",   AMBER,    "react-hot-toast with custom dark-themed style for success / error feedback."),
    ("⚡", "Skeleton Loaders",      CYAN,     "Shimmer placeholder cards shown while API requests are in-flight — no blank flash."),
    ("🔍", "Live Search",           PRIMARY,  "Search debounce via useCallback + useEffect — triggers new API fetch on each change."),
    ("🏷️", "Filter Chips",          PURPLE2,  "Category & type chips with active state (colour changes) — instant visual feedback."),
]

for i, (icon, title, col, desc) in enumerate(design_points):
    col_i = i // 4
    row_i = i %  4
    dx = Inches(0.5 + col_i * 6.45)
    dy = Inches(1.45) + row_i * Inches(1.45)
    add_rect(sl, dx, dy, Inches(6.2), Inches(1.3), fill=SURFACE2, line=col, line_w=Pt(0.8))
    add_text(sl, icon, dx + Inches(0.15), dy + Inches(0.18), Inches(0.5), Inches(0.5), size=22)
    add_text(sl, title, dx + Inches(0.75), dy + Inches(0.18), Inches(5.4), Inches(0.4), size=14, bold=True, color=col)
    add_text(sl, desc, dx + Inches(0.15), dy + Inches(0.65), Inches(5.9), Inches(0.58), size=11, color=MUTED)

accent_bar(sl)
slide_number(sl, 10)

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 11 — FOLDER STRUCTURE
# ══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
slide_bg(sl)

add_text(sl, "Project Folder Structure", Inches(0.6), Inches(0.5), Inches(10), Inches(0.6),
         size=34, bold=True, color=WHITE, name="Calibri")
add_rect(sl, Inches(0.6), Inches(1.15), Inches(1.8), Inches(0.04), fill=PRIMARY)

tree_left = [
    ("fsd final project/",        WHITE,   0, PRIMARY),
    ("├── package.json",          MUTED,   1, None),
    ("├── .gitignore",            MUTED,   1, None),
    ("├── server/",               WHITE,   1, CYAN),
    ("│   ├── server.js",         MUTED,   2, None),
    ("│   ├── .env",              ROSE,    2, None),
    ("│   ├── models/",           WHITE,   2, GREEN),
    ("│   │   ├── News.js",       MUTED,   3, None),
    ("│   │   ├── LostFound.js",  MUTED,   3, None),
    ("│   │   └── Textbook.js",   MUTED,   3, None),
    ("│   └── routes/",           WHITE,   2, GREEN),
    ("│       ├── news.js",       MUTED,   3, None),
    ("│       ├── lostfound.js",  MUTED,   3, None),
    ("│       └── textbooks.js",  MUTED,   3, None),
]
tree_right = [
    ("└── client/",               WHITE,   1, CYAN),
    ("    ├── index.html",        MUTED,   2, None),
    ("    ├── vite.config.js",    MUTED,   2, None),
    ("    ├── package.json",      MUTED,   2, None),
    ("    └── src/",              WHITE,   2, PRIMARY),
    ("        ├── main.jsx",      MUTED,   3, None),
    ("        ├── App.jsx",       MUTED,   3, None),
    ("        ├── index.css",     MUTED,   3, None),
    ("        ├── api/index.js",  MUTED,   3, None),
    ("        ├── components/",   WHITE,   3, CYAN),
    ("        │   ├── Navbar.jsx",MUTED,   4, None),
    ("        │   ├── Footer.jsx",MUTED,   4, None),
    ("        │   └── Modal.jsx", MUTED,   4, None),
    ("        └── pages/",        WHITE,   3, CYAN),
    ("            ├── Home.jsx",  MUTED,   4, None),
    ("            ├── NewsPage.jsx",    MUTED, 4, None),
    ("            ├── LostFoundPage.jsx",MUTED,4, None),
    ("            └── TextbooksPage.jsx",MUTED,4, None),
]

# Code block background
add_rect(sl, Inches(0.5), Inches(1.45), Inches(12.33), Inches(5.7), fill=SURFACE2)

for i, (text, col, indent, hl) in enumerate(tree_left):
    ty = Inches(1.55) + i * Inches(0.36)
    add_text(sl, text, Inches(0.7 + indent*0.18), ty, Inches(5.8), Inches(0.33),
             size=11, color=col, name="Courier New")
    if hl:
        add_rect(sl, Inches(0.55), ty, Inches(0.04), Inches(0.3), fill=hl)

for i, (text, col, indent, hl) in enumerate(tree_right):
    ty = Inches(1.55) + i * Inches(0.36)
    add_text(sl, text, Inches(6.7 + indent*0.18), ty, Inches(6.2), Inches(0.33),
             size=11, color=col, name="Courier New")
    if hl:
        add_rect(sl, Inches(6.55), ty, Inches(0.04), Inches(0.3), fill=hl)

accent_bar(sl)
slide_number(sl, 11)

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 12 — FUTURE ENHANCEMENTS & CONCLUSION
# ══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
slide_bg(sl)
glow_orb(sl, Inches(1),    Inches(6),   3.5, PRIMARY)
glow_orb(sl, Inches(12.5), Inches(1.5), 3.0, CYAN)

add_text(sl, "Future Enhancements & Conclusion", Inches(0.6), Inches(0.48), Inches(12), Inches(0.6),
         size=30, bold=True, color=WHITE, name="Calibri")
add_rect(sl, Inches(0.6), Inches(1.12), Inches(2.2), Inches(0.04), fill=CYAN)

# Future enhancements
enhancements = [
    ("🔐", "User Authentication",    "JWT-based login / register — profile pages with posting history"),
    ("📸", "Image Uploads",          "Multer middleware for photo uploads on news & L&F items"),
    ("💬", "In-app Messaging",       "Real-time chat between students using Socket.io"),
    ("🔎", "Advanced Search",        "Elasticsearch integration for fuzzy & phonetic matching"),
    ("📈", "Admin Dashboard",        "Analytics view — daily posts, active items, textbook turnover"),
    ("📲", "PWA Support",            "Install as mobile app — offline mode & push notifications"),
]

for i, (icon, title, desc) in enumerate(enhancements):
    ex = Inches(0.5 + (i % 3) * 4.18)
    ey = Inches(1.45) + (i // 3) * Inches(1.5)
    add_rect(sl, ex, ey, Inches(4.0), Inches(1.32), fill=SURFACE2, line=CYAN, line_w=Pt(0.7))
    add_text(sl, icon,  ex + Inches(0.15), ey + Inches(0.15), Inches(0.5),  Inches(0.5), size=20)
    add_text(sl, title, ex + Inches(0.72), ey + Inches(0.17), Inches(3.1),  Inches(0.4), size=13, bold=True, color=WHITE)
    add_text(sl, desc,  ex + Inches(0.15), ey + Inches(0.65), Inches(3.75), Inches(0.58), size=11, color=MUTED)

# Conclusion bar
add_rect(sl, Inches(0.5), Inches(5.75), Inches(12.33), Inches(1.35), fill=PRIMARY)
add_text(sl,
    "CampusHub demonstrates a production-ready MERN stack architecture — "
    "from MongoDB schemas and Express REST APIs to a React frontend with premium dark UI, "
    "live search, category filtering, and direct student-to-student contact.",
    Inches(0.75), Inches(5.9), Inches(11.9), Inches(1.1),
    size=13, color=WHITE, align=PP_ALIGN.CENTER, wrap=True)

accent_bar(sl, CYAN)
slide_number(sl, 12)


out = r"c:\Users\admin\Documents\jairajcj_projects\fsd final project\CampusHub_Presentation.pptx"
prs.save(out)
print(f"[OK] Saved: {out}")
